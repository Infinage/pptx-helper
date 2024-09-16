import pptx
import os
import shutil

class PPTHelper:

    # 1 inch in EMU divided by 96 dpi
    EMU_TO_PX: float = 914400 / 96

    @staticmethod
    def emu_to_px(units: int) -> int:
        """
        Convert EMU (English Metric Units) to pixels.
        """
        return int(units / PPTHelper.EMU_TO_PX)

    def __init__(self, input_fpath: str, output_dpath: str = ".") -> None:
        self.input_fpath = input_fpath
        self.output_dpath = output_dpath

        # If not exists, create empty path
        os.makedirs(self.output_dpath, exist_ok=True)

        self.presentation = pptx.Presentation(self.input_fpath)

        # Get slide dimensions (in pixels)
        self.slide_width = PPTHelper.emu_to_px(self.presentation.slide_width or 0)
        self.slide_height = PPTHelper.emu_to_px(self.presentation.slide_height or 0)

    def to_html(self) -> None:

        # Copy style.css
        shutil.copy2("style.css", os.path.join(self.output_dpath, "style.css"))


        # Index.html
        with open(os.path.join(self.output_dpath, "index.html"), "w") as f:
            f.write(f"""
                <html>
                <head>
                    <link rel='stylesheet' href='style.css'>
                </head>
                <div>
                    <p>Input File Path: {self.input_fpath} </p>
                </div>
                <ol class='ppt-slides'>
                {"".join("<li><a href='slide_" + str(slide_id) + ".html'>Slide #" + str(slide_id + 1) + "</a>" for slide_id in range(len(self.presentation.slides)))}
                </ol>
                </html>
        """)

        # Slides_x.html
        for slide_id, slide in enumerate(self.presentation.slides):
            markup = self.generate_slide_overview(slide_id)
            with open(os.path.join(self.output_dpath, f"slide_{slide_id}.html"), "w") as f:
                f.write(markup)

            # Generate fragment for each shape type - frag_x_y.html
            for shape_id, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    markup = self.text_frame_to_fragment(slide_id, shape_id)

                elif shape.has_table:
                    markup = self.table_to_fragment(slide_id, shape_id)

                    # Generate fragment for each cell's textframe
                    for row_num, row in enumerate(shape.table.rows):
                        for col_num, cell in enumerate(row.cells):
                            cell_fragment = self.cell_text_frame_to_fragment(slide_id, shape_id, row_num, col_num)
                            with open(os.path.join(self.output_dpath, f"frag_{slide_id}_{shape_id}_({row_num},{col_num}).html"), "w") as f:
                                f.write(f"<html><head><link rel='stylesheet' href='style.css'></head>\n{cell_fragment}</html>")

                else:
                    markup = f"Sorry, shape type: '{shape.shape_type}' not supported yet."

                with open(os.path.join(self.output_dpath, f"frag_{slide_id}_{shape_id}.html"), "w") as f:
                    f.write(f"<html><head><link rel='stylesheet' href='style.css'></head>\n{markup}</html>")
                    #f.write(f"{markup}")

    def generate_slide_overview(self, slide_id: int) -> str:
        slide_html = (
                f'<html><head><link rel="stylesheet" href="style.css"></head>\n'
                f'<div class="slide" style="width:{self.slide_width}px; height:{self.slide_height}px;">\n'
        )

        # Loop through each shape in the slide
        for shape_id, shape in enumerate(self.presentation.slides[slide_id].shapes):
            # Get shape position and dimensions in pixels
            left = PPTHelper.emu_to_px(shape.left)
            top = PPTHelper.emu_to_px(shape.top)
            width = PPTHelper.emu_to_px(shape.width)
            height = PPTHelper.emu_to_px(shape.height)

            # Create a rectangle for each shape
            shape_html = (
                f"<a class='nav-link' href='frag_{slide_id}_{shape_id}.html'>"
                f'<div class="shape" style="left:{left}px; top:{top}px; width:{width}px; height:{height}px;">\n'
                f'<p>Shape ID: {shape_id}<br>Type: {shape.shape_type}</p>\n'
                '</div></a>\n'
            )

            slide_html += shape_html

        slide_html += (
            f'<div class="slide-footer">'
                f'<a href="index.html" class="nav-link">'
                    f'<p class="slide-id">Slide ID: {slide_id}</p>'
                f'</a>'
            f'</div></html>'
        )
        return slide_html

    def text_frame_to_fragment(self, slide_id: int, shape_id: int) -> str:
        frag = (
            f"<div class='text-frame'><p class='slide-shape-id'>"
                f"<a href='index.html' class='nav-link'>Slide ID: {slide_id}</a> / "
                f"<a href='slide_{slide_id}.html' class='nav-link'>Shape ID: {shape_id}</a>"
            f"</p>"
        )
        for para_id, para in enumerate(self.presentation.slides[slide_id].shapes[shape_id].text_frame.paragraphs):
            frag += f"<div class='text-frame-para'><p class='para-id'>Para ID: {para_id}</p>"
            for run_id, run in enumerate(para.runs):
                frag += f"<div class='text-frame-run'><p class='run-id'>Run ID: {run_id}</p><p class='run-text'>{run.text}</p></div>"
            frag += "</div>"
        frag += "</div>"
        return frag

    def cell_text_frame_to_fragment(self, slide_id: int, shape_id: int, row_num: int, col_num: int) -> str:
        frag = (
            f"<div class='text-frame'>"
                "<p class='slide-shape-id'>"
                    f"<a href='index.html' class='nav-link'>Slide ID: {slide_id}</a> / "
                    f"<a href='slide_{slide_id}.html' class='nav-link'>Shape ID: {shape_id}</a>"
                "</p>"
                "<p class='slide-shape-id'>"
                    f"<a href='frag_{slide_id}_{shape_id}.html' class='nav-link'>Row / Col: {shape_id}, {col_num}</a>"
                "</p>"
            f"</p>"
        )
        for para_id, para in enumerate(self.presentation.slides[slide_id].shapes[shape_id].table.cell(row_num, col_num).text_frame.paragraphs):
            frag += f"<div class='text-frame-para'><p class='para-id'>Para ID: {para_id}</p>"
            for run_id, run in enumerate(para.runs):
                frag += f"<div class='text-frame-run'><p class='run-id'>Run ID: {run_id}</p><p class='run-text'>{run.text}</p></div>"
            frag += "</div>"
        frag += "</div>"
        return frag

    def table_to_fragment(self, slide_id: int, shape_id: int) -> str:
        table = self.presentation.slides[slide_id].shapes[shape_id].table
        frag = (
            f"<div class='table-div'>"
                "<p class='slide-shape-id'>"
                    f"<a href='index.html' class='nav-link'>Slide ID: {slide_id}</a> / "
                    f"<a href='slide_{slide_id}.html' class='nav-link'>Shape ID: {shape_id}</a>"
                f"</p>"
            f"<table>"
        )

        for row_num, row in enumerate(table.rows):
            frag += "<tr>"
            for col_num, cell in enumerate(row.cells):
                if not cell.is_spanned:
                    frag += (
                        f"<td colspan='{cell.span_width}' rowspan='{cell.span_height}'>"
                            f"<a class='nav-link' href='frag_{slide_id}_{shape_id}_({row_num},{col_num}).html'>"
                            f"<div class='table-cell-text'>{cell.text}</div>"
                        "</td>"
                    )
            frag += "</tr>"

        frag += "</table></div>"
        return frag
